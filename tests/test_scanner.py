import os
import time
import tempfile
import json
import threading
from pathlib import Path
from unittest.mock import MagicMock, patch
from PIL import Image
import pytest
from backend.scanner import (
    get_system_ram, get_auto_config,
    PII_PATTERNS, VENDOR_PATTERNS, calculate_shannon_entropy, detect_secrets, mask_secret,
    parse_ai_response,
    get_file_text_content, inspect_text,
    verify_text_file_with_ai, inspect_image, process_heic_image,
    get_scannable_files, get_optimized_image_path,
    get_ollama_address, get_client, ensure_ollama_running,
    run_scan, start_scan_thread, stop_scan,
    ScannerState, scan_state, locate_text_pii_matches,
    get_inference_response, get_model_provider,
    get_active_vision_model, get_active_text_model,
    calculate_file_checksum
)
from backend.create_icon import generate_argus_icon
from backend.state import save_cache, load_cache

def test_generate_argus_icon(tmp_path, monkeypatch):
    # Verify icon generator runs and produces valid image assets
    generate_argus_icon()
    assets_dir = Path(__file__).parent.parent / "frontend" / "assets"
    assert (assets_dir / "argus-icon.ico").exists()
    assert (assets_dir / "argus-icon.png").exists()

def test_get_auto_config():
    # < 8GB
    cfg_low = get_auto_config(4 * (1024 ** 3))
    assert cfg_low["concurrency"] == 1
    assert cfg_low["image_optimization"] == "low"

    # 8GB <= RAM < 16GB
    cfg_med = get_auto_config(12 * (1024 ** 3))
    assert cfg_med["concurrency"] == 2
    assert cfg_med["image_optimization"] == "medium"

    # >= 16GB
    cfg_high = get_auto_config(32 * (1024 ** 3))
    assert cfg_high["concurrency"] == 4
    assert cfg_high["image_optimization"] == "medium"

def test_get_system_ram():
    ram = get_system_ram()
    assert isinstance(ram, (int, float))
    assert ram > 0

def test_get_system_ram_exception():
    with patch("sys.platform", "darwin"):
        with patch("subprocess.check_output", side_effect=Exception("sysctl failed")):
            fallback_ram = get_system_ram()
            assert fallback_ram == 8 * 1024 * 1024 * 1024

def test_get_system_ram_darwin():
    with patch("sys.platform", "darwin"):
        with patch("subprocess.check_output", return_value=b"17179869184\n"):
            ram = get_system_ram()
            assert ram == 17179869184

def test_get_system_ram_linux():
    from unittest.mock import mock_open
    fake_meminfo = "MemTotal:       32890624 kB\nMemFree: 1000 kB\n"
    with patch("sys.platform", "linux"):
        with patch("builtins.open", mock_open(read_data=fake_meminfo)):
            ram = get_system_ram()
            assert ram == 32890624 * 1024


def test_pii_regex_patterns():
    ssn_text = "The user SSN is 123-45-6789 and needs protection."
    assert PII_PATTERNS["SSN"].search(ssn_text) is not None

    cc_text = "Card number 4111-2222-3333-4444 should be blocked."
    assert PII_PATTERNS["Credit Card"].search(cc_text) is not None

    pk_text = "-----BEGIN RSA PRIVATE KEY-----\nMIIEowIBAAKCAQEA0...\n-----END RSA PRIVATE KEY-----"
    assert VENDOR_PATTERNS["Private Key"].search(pk_text) is not None

    api_text = 'const apiKey = "sk-proj-1234567890abcdef1234567890abcdef12345678";'
    assert VENDOR_PATTERNS["OpenAI API Key"].search(api_text) is not None

    clean_text = "Just a regular text document with no sensitive data."
    for pattern in PII_PATTERNS.values():
        assert pattern.search(clean_text) is None
    for pattern in VENDOR_PATTERNS.values():
        assert pattern.search(clean_text) is None

def test_detect_secrets():
    content = """Line 1: Clean line
Line 2: AWS Key is AKIA1234567890ABCDEF here
Line 3: Secret token is "a1b2c3d4e5f6g7h8i9j0k1l2m3n4o5p6"
"""
    # Test Tier 1 (AWS)
    matches = detect_secrets(content)
    aws_matches = [m for m in matches if m["pattern_name"] == "AWS Access Key"]
    assert len(aws_matches) == 1
    assert aws_matches[0]["line_number"] == 2
    assert "..." in aws_matches[0]["match_text"]  # Should be masked
    
    # Test Tier 2 & 3 (Generic Token + Context)
    gen_matches = [m for m in matches if m["pattern_name"] == "Generic API Token"]
    assert len(gen_matches) == 1
    assert gen_matches[0]["line_number"] == 3
    assert "..." in gen_matches[0]["match_text"]

def test_entropy_and_masking():
    # High entropy string
    assert calculate_shannon_entropy("a1b2c3d4e5f6g7h8") > 3.0
    # Low entropy string
    assert calculate_shannon_entropy("aaaaaaaaaaaaaaaa") == 0.0
    
    # Masking
    assert mask_secret("short") == "***"
    assert mask_secret("AKIA1234567890ABCDEF") == "AKIA1234...CDEF"

def test_parse_ai_response():
    valid_json = '{"compromised": true, "reason": "Contains SSN", "items": [{"label": "SSN Card", "box_2d": [100, 200, 300, 400], "description": "SSN"}], "snippets": ["123-45-6789"]}'
    parsed = parse_ai_response(valid_json)
    assert parsed["compromised"] is True
    assert parsed["reason"] == "Contains SSN"
    assert len(parsed["items"]) == 1
    assert parsed["items"][0]["box_2d"] == [100.0, 200.0, 300.0, 400.0]
    assert parsed["snippets"] == ["123-45-6789"]

    # Invalid box coordinates handled gracefully
    bad_box_json = '{"compromised": true, "items": [{"label": "Card", "box_2d": "invalid"}]}'
    parsed_bad = parse_ai_response(bad_box_json)
    assert parsed_bad["items"][0]["box_2d"] is None

    markdown_json = '```json\n{"compromised": false, "reason": "Clean", "items": [], "snippets": []}\n```'
    parsed = parse_ai_response(markdown_json)
    assert parsed["compromised"] is False

    # Non-dict JSON string
    assert parse_ai_response("[1, 2, 3]")["compromised"] is False

    # Invalid items and snippets types
    non_list_json = '{"compromised": true, "items": "not_a_list", "snippets": "not_a_list"}'
    parsed_non_list = parse_ai_response(non_list_json)
    assert parsed_non_list["compromised"] is True
    assert parsed_non_list["items"] == []
    assert parsed_non_list["snippets"] == []

    # Regex extracted JSON from conversational wrapper
    convo_json = 'Sure! Here is the DLP result:\n{"compromised": true, "reason": "Found driver license", "items": []}\nHope this helps!'
    parsed_convo = parse_ai_response(convo_json)
    assert parsed_convo["compromised"] is True
    assert parsed_convo["reason"] == "Found driver license"

    # Conversational JSON with non-dict regex match
    convo_bad = 'Result is: {"key": [1, 2, 3]}'
    assert parse_ai_response(convo_bad)["compromised"] is False

    fallback_compromised = "Result: compromised = true, contains passport"
    parsed_comp = parse_ai_response(fallback_compromised)
    assert parsed_comp["compromised"] is True

    fallback_clean = "File appears clean and safe."
    parsed_clean = parse_ai_response(fallback_clean)
    assert parsed_clean["compromised"] is False

def test_image_optimization_downsampling(tmp_path):
    # Large image that needs downsampling
    large_img_path = tmp_path / "large_photo.png"
    large_img = Image.new("RGBA", (2000, 2000), color="yellow")
    large_img.save(str(large_img_path))

    opt_path, is_temp = get_optimized_image_path(str(large_img_path), "medium")
    assert is_temp is True
    assert os.path.exists(opt_path)
    with Image.open(opt_path) as img:
        assert max(img.size) <= 1024
    if os.path.exists(opt_path):
        os.remove(opt_path)

    # Original mode (no downsampling)
    orig_path, orig_temp = get_optimized_image_path(str(large_img_path), "original")
    assert orig_temp is False
    assert orig_path == str(large_img_path)

    # Small image (no downsampling needed)
    small_path = tmp_path / "small.jpg"
    small_img = Image.new("RGB", (200, 200), color="red")
    small_img.save(str(small_path))
    res_path, res_temp = get_optimized_image_path(str(small_path), "medium")
    assert res_temp is False
    assert res_path == str(small_path)

def test_locate_text_pii_matches():
    # Empty content
    assert locate_text_pii_matches("") == []
    assert locate_text_pii_matches(None) == []

    content = """Line 1: Clean line
Line 2: User SSN is 123-45-6789 here
Line 3: Clean line
Line 4: Secret apiKey = "sk-proj-1234567890abcdef1234567890abcdef12345678"
Line 5: Custom sensitive token secret_token_xyz
"""
    # Test regex only
    matches = locate_text_pii_matches(content)
    assert len(matches) >= 2
    
    ssn_matches = [m for m in matches if m["pattern_name"] == "SSN"]
    assert len(ssn_matches) > 0
    assert ssn_matches[0]["line_number"] == 2
    
    openai_matches = [m for m in matches if m["pattern_name"] == "OpenAI API Key"]
    assert len(openai_matches) > 0
    assert openai_matches[0]["line_number"] == 4

    # Test with AI snippets
    matches_with_ai = locate_text_pii_matches(content, ai_snippets=["secret_token_xyz", "", None])
    ai_matches = [m for m in matches_with_ai if m["source"] == "ai"]
    assert len(ai_matches) == 1
    assert ai_matches[0]["line_number"] == 5
    assert ai_matches[0]["match_text"] == "secret_token_xyz"

def test_get_ollama_address_and_client(monkeypatch):
    import backend.scanner as scanner
    monkeypatch.setattr(scanner, "load_settings", lambda: {"ollama_address": "http://10.0.0.1:11434/"})
    addr = get_ollama_address()
    assert addr == "http://10.0.0.1:11434"
    client = get_client()
    assert client is not None

def test_ensure_ollama_running(monkeypatch):
    # Case 1: Already running
    with patch("urllib.request.urlopen") as mock_url:
        mock_url.return_value = True
        success, msg = ensure_ollama_running()
        assert success is True
        assert "running" in msg

    # Case 2: Remote server unreachable
    with patch("urllib.request.urlopen", side_effect=Exception("Connection refused")):
        monkeypatch.setattr("backend.scanner.get_ollama_address", lambda: "http://remote-server:11434")
        success, msg = ensure_ollama_running()
        assert success is False
        assert "Could not connect to remote Ollama server" in msg

    # Case 3: Local server launched successfully
    with patch("urllib.request.urlopen", side_effect=[Exception("Refused"), True]):
        monkeypatch.setattr("backend.scanner.get_ollama_address", lambda: "http://127.0.0.1:11434")
        with patch("subprocess.Popen") as mock_popen:
            success, msg = ensure_ollama_running()
            assert success is True

    # Case 4: FileNotFoundError
    with patch("urllib.request.urlopen", side_effect=Exception("Refused")):
        monkeypatch.setattr("backend.scanner.get_ollama_address", lambda: "http://127.0.0.1:11434")
        with patch("subprocess.Popen", side_effect=FileNotFoundError("Ollama not found")):
            s_fnf, m_fnf = ensure_ollama_running()
            assert s_fnf is False
            assert "executable not found" in m_fnf

    # Case 5: Unexpected error
    with patch("urllib.request.urlopen", side_effect=Exception("Refused")):
        monkeypatch.setattr("backend.scanner.get_ollama_address", lambda: "http://127.0.0.1:11434")
        with patch("subprocess.Popen", side_effect=Exception("Fatal Error")):
            s_err, m_err = ensure_ollama_running()
            assert s_err is False
            assert "Unexpected error" in m_err

def test_get_file_text_content_all_formats(tmp_path):
    # Plain text
    txt_file = tmp_path / "sample.txt"
    txt_file.write_text("Plain text sample", encoding="utf-8")
    assert get_file_text_content(txt_file) == "Plain text sample"

    # Word DOCX
    import docx
    doc_path = tmp_path / "sample.docx"
    doc = docx.Document()
    doc.add_paragraph("DOCX Paragraph Content")
    t = doc.add_table(rows=1, cols=1)
    t.rows[0].cells[0].text = "TableCell"
    doc.save(str(doc_path))
    docx_text = get_file_text_content(doc_path)
    assert "DOCX Paragraph Content" in docx_text
    assert "TableCell" in docx_text

    # Excel XLSX
    import openpyxl
    xlsx_path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Account Number"
    ws["B1"] = "987654321"
    wb.save(str(xlsx_path))
    wb.close()
    xlsx_text = get_file_text_content(xlsx_path)
    assert "Account Number" in xlsx_text
    assert "987654321" in xlsx_text

    # PowerPoint PPTX
    import pptx
    pptx_path = tmp_path / "sample.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Presentation Secret"
    prs.save(str(pptx_path))
    pptx_text = get_file_text_content(pptx_path)
    assert "Presentation Secret" in pptx_text

    # PDF
    import pypdf
    pdf_path = tmp_path / "sample.pdf"
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with open(pdf_path, "wb") as f:
        writer.write(f)
    assert isinstance(get_file_text_content(pdf_path), str)

    # Corrupt document handling
    bad_pdf = tmp_path / "corrupt.pdf"
    bad_pdf.write_text("invalid binary")
    assert get_file_text_content(bad_pdf) == ""

    bad_docx = tmp_path / "corrupt.docx"
    bad_docx.write_text("invalid binary")
    assert get_file_text_content(bad_docx) == ""

    bad_xlsx = tmp_path / "corrupt.xlsx"
    bad_xlsx.write_text("invalid binary")
    assert get_file_text_content(bad_xlsx) == ""

    bad_pptx = tmp_path / "corrupt.pptx"
    bad_pptx.write_text("invalid binary")
    assert get_file_text_content(bad_pptx) == ""

def test_inspect_text_and_verify_with_ai(tmp_path):
    # Regex only mode
    ssn_file = tmp_path / "ssn.txt"
    ssn_file.write_text("111-22-3333", encoding="utf-8")
    r_only = inspect_text(ssn_file, "regex_only")
    assert r_only["compromised"] is True

    # LLM Only with Mock
    mock_chat_response = {"message": {"content": '{"compromised": true, "reason": "Found bank info"}'}}
    with patch("backend.scanner.get_client") as mock_client_fn:
        mock_client = MagicMock()
        mock_client.chat.return_value = mock_chat_response
        mock_client_fn.return_value = mock_client

        clean_file = tmp_path / "statement.txt"
        clean_file.write_text("Bank statement for John", encoding="utf-8")
        
        r_llm = inspect_text(clean_file, "llm_only")
        assert r_llm["compromised"] is True
        assert r_llm["reason"] == "Found bank info"

        # Test verify_text_file_with_ai
        verified = verify_text_file_with_ai(str(clean_file))
        assert verified["compromised"] is True

        # Test exception handling in inspect_text and verify_text_file_with_ai
        mock_client.chat.side_effect = Exception("Ollama disconnected")
        err_res = inspect_text(clean_file, "llm_only")
        assert err_res["compromised"] is False
        assert "Error inspecting text" in err_res["reason"]

        err_verify = verify_text_file_with_ai(str(clean_file))
        assert err_verify["compromised"] is False

def test_inspect_image_and_heic(tmp_path):
    img_path = tmp_path / "photo.jpg"
    img = Image.new("RGB", (100, 100), color="green")
    img.save(str(img_path))

    mock_chat_response = {"message": {"content": '{"compromised": false, "reason": "No PII"}'}}
    with patch("backend.scanner.get_client") as mock_client_fn:
        mock_client = MagicMock()
        mock_client.chat.return_value = mock_chat_response
        mock_client_fn.return_value = mock_client

        result = inspect_image(str(img_path), "medium")
        assert result["compromised"] is False

        # Test HEIC process
        heic_result = process_heic_image(img_path, "medium")
        assert heic_result["compromised"] is False

        # Test inspect_image exception
        mock_client.chat.side_effect = Exception("Model failed")
        err_img = inspect_image(str(img_path), "medium")
        assert err_img["compromised"] is False

def test_run_scan_lifecycle(tmp_path):
    folder = tmp_path / "scandir"
    folder.mkdir()
    
    file1 = folder / "test_ssn.txt"
    file1.write_text("999-88-7777 secret ssn", encoding="utf-8")
    
    file2 = folder / "clean.txt"
    file2.write_text("totally clean document", encoding="utf-8")

    results_received = []
    def save_cb(flagged):
        results_received.extend(flagged)

    # Run scan with regex_only mode and auto_delete=True
    with patch("backend.scanner.load_settings", lambda: {
        "concurrency": "2",
        "image_optimization": "medium",
        "text_scan_mode": "regex_only",
        "auto_delete": True
    }):
        run_scan([str(folder)], save_cb, rescan_all=True)
        assert len(results_received) >= 1
        assert any("test_ssn.txt" in f["file"] for f in results_received)
        # auto_deleted should be True for regex_only match
        assert results_received[0]["auto_deleted"] is True
        assert results_received[0]["compromised"] is True
        assert not file1.exists()

def test_run_scan_cache_compromised_reload(tmp_path):
    folder = tmp_path / "cached_dir"
    folder.mkdir()
    
    flagged_file = folder / "preflagged.txt"
    flagged_file.write_text("123-45-6789 confidential")
    
    mtime = os.path.getmtime(str(flagged_file))
    save_cache({
        str(flagged_file): {
            "mtime": mtime,
            "result": {"compromised": True, "reason": "Pre-flagged in cache", "needs_ai_verification": True}
        }
    })

    results = []
    with patch("backend.scanner.load_settings", lambda: {"concurrency": "1", "image_optimization": "medium", "text_scan_mode": "regex_only"}):
        run_scan([str(folder)], lambda r: results.extend(r), rescan_all=False)
        assert len(results) == 1
        assert results[0]["file"] == str(flagged_file)
        assert results[0]["reason"] == "Pre-flagged in cache"
        assert results[0]["compromised"] is True

def test_calculate_file_checksum(tmp_path):
    # Valid file with known content
    test_file = tmp_path / "hello.txt"
    test_file.write_text("Hello World", encoding="utf-8")
    chk = calculate_file_checksum(test_file)
    import hashlib
    expected_hash = hashlib.sha256(b"Hello World").hexdigest()
    assert chk == expected_hash

    # Empty file
    empty_file = tmp_path / "empty.txt"
    empty_file.write_text("", encoding="utf-8")
    assert calculate_file_checksum(empty_file) == hashlib.sha256(b"").hexdigest()

    # Non-existent file returns empty string
    assert calculate_file_checksum(tmp_path / "does_not_exist.txt") == ""

def test_run_scan_altered_file_detection(tmp_path):
    folder = tmp_path / "scan_alter_dir"
    folder.mkdir()
    target_file = folder / "document.txt"

    # Step 1: Initial clean file
    target_file.write_text("Clean initial document with no sensitive data", encoding="utf-8")
    results1 = []
    with patch("backend.scanner.load_settings", lambda: {"concurrency": "1", "image_optimization": "medium", "text_scan_mode": "regex_only"}):
        run_scan([str(folder)], lambda r: results1.extend(r), rescan_all=False)
        assert len(results1) == 0
        cache1 = load_cache()
        assert str(target_file) in cache1
        assert cache1[str(target_file)]["checksum"] == calculate_file_checksum(target_file)
        assert cache1[str(target_file)]["result"]["compromised"] is False

    # Step 2: Alter the file by adding an SSN
    target_file.write_text("Document updated: SSN is 123-45-6789 confidential", encoding="utf-8")
    results2 = []
    with patch("backend.scanner.load_settings", lambda: {"concurrency": "1", "image_optimization": "medium", "text_scan_mode": "regex_only"}):
        run_scan([str(folder)], lambda r: results2.extend(r), rescan_all=False)
        # Checksum change was detected, file was re-scanned and flagged
        assert len(results2) == 1
        assert results2[0]["file"] == str(target_file)
        assert results2[0]["compromised"] is True
        cache2 = load_cache()
        assert cache2[str(target_file)]["checksum"] == calculate_file_checksum(target_file)
        assert cache2[str(target_file)]["result"]["compromised"] is True

    # Step 3: Alter the file again to remediate / clean it
    target_file.write_text("Remediated document: all sensitive data removed", encoding="utf-8")
    results3 = []
    with patch("backend.scanner.load_settings", lambda: {"concurrency": "1", "image_optimization": "medium", "text_scan_mode": "regex_only"}):
        run_scan([str(folder)], lambda r: results3.extend(r), rescan_all=False)
        # Checksum change was detected, file was re-scanned and cleared
        assert len(results3) == 0
        cache3 = load_cache()
        assert cache3[str(target_file)]["checksum"] == calculate_file_checksum(target_file)
        assert cache3[str(target_file)]["result"]["compromised"] is False

def test_run_scan_unmodified_file_skips_inspection(tmp_path):
    folder = tmp_path / "skip_dir"
    folder.mkdir()
    clean_file = folder / "clean.txt"
    clean_file.write_text("Static clean content", encoding="utf-8")

    # Initial scan to populate cache
    with patch("backend.scanner.load_settings", lambda: {"concurrency": "1", "image_optimization": "medium", "text_scan_mode": "regex_only"}):
        run_scan([str(folder)], lambda r: None, rescan_all=False)

    # Second scan: mock inspect_text to ensure it is NOT invoked
    with patch("backend.scanner.inspect_text") as mock_inspect:
        with patch("backend.scanner.load_settings", lambda: {"concurrency": "1", "image_optimization": "medium", "text_scan_mode": "regex_only"}):
            run_scan([str(folder)], lambda r: None, rescan_all=False)
            mock_inspect.assert_not_called()

def test_run_scan_legacy_cache_migration(tmp_path):
    folder = tmp_path / "legacy_dir"
    folder.mkdir()
    doc_file = folder / "legacy_doc.txt"
    doc_file.write_text("Legacy doc content", encoding="utf-8")

    mtime = os.path.getmtime(str(doc_file))
    # Cache without 'checksum' field (legacy format)
    save_cache({
        str(doc_file): {
            "mtime": mtime,
            "result": {"compromised": False, "reason": "No PII"}
        }
    })

    with patch("backend.scanner.load_settings", lambda: {"concurrency": "1", "image_optimization": "medium", "text_scan_mode": "regex_only"}):
        run_scan([str(folder)], lambda r: None, rescan_all=False)
        cache = load_cache()
        # Legacy entry was migrated with checksum and size
        assert "checksum" in cache[str(doc_file)]
        assert cache[str(doc_file)]["checksum"] == calculate_file_checksum(doc_file)
        assert "size" in cache[str(doc_file)]

def test_start_and_stop_scan(tmp_path):
    scan_state.reset()
    folder = tmp_path / "empty_dir"
    folder.mkdir()

    hold_event = threading.Event()
    release_event = threading.Event()

    def mock_run_scan(*args, **kwargs):
        scan_state.is_scanning = True
        hold_event.set()
        release_event.wait(timeout=2.0)
        scan_state.is_scanning = False

    with patch("backend.scanner.run_scan", side_effect=mock_run_scan):
        assert start_scan_thread([str(folder)], lambda r: None) is True
        assert hold_event.wait(timeout=1.0) is True
        # Second start while running returns False
        assert start_scan_thread([str(folder)], lambda r: None) is False
        
        # Abort scan
        stop_scan()
        assert scan_state.should_stop is True
        release_event.set()

    # Wait for thread to finish and test starting again after stop
    time.sleep(0.1)
    with patch("backend.scanner.run_scan", side_effect=lambda *a, **k: None):
        assert start_scan_thread([str(folder)], lambda r: None) is True
        stop_scan(timeout=0.5)



# ============================================================================
# INFERENCE CLIENT DISPATCH TESTS
# ============================================================================

def test_get_model_provider_default():
    provider = get_model_provider()
    assert provider == "ollama"

def test_get_model_provider_from_settings(tmp_path):
    from backend.state import save_settings
    save_settings({"model_provider": "local_gguf"})
    assert get_model_provider() == "local_gguf"

def test_get_active_vision_model_default():
    model = get_active_vision_model()
    assert model == "gemma4:12b"

def test_get_active_text_model_default():
    model = get_active_text_model()
    assert model == "gemma4:12b"

def test_get_active_vision_model_custom(tmp_path):
    from backend.state import save_settings
    save_settings({"vision_model_name": "llava:7b"})
    assert get_active_vision_model() == "llava:7b"

def test_get_active_text_model_custom(tmp_path):
    from backend.state import save_settings
    save_settings({"text_model_name": "mistral:7b"})
    assert get_active_text_model() == "mistral:7b"

def test_get_inference_response_ollama():
    mock_client = MagicMock()
    mock_client.chat.return_value = {"message": {"content": '{"compromised": false}'}}
    with patch("backend.scanner.get_model_provider", return_value="ollama"):
        with patch("backend.scanner.get_client", return_value=mock_client):
            result = get_inference_response(
                messages=[{"role": "user", "content": "test"}],
                model_name="gemma4:12b",
                options={"temperature": 0.0}
            )
            assert result["message"]["content"] == '{"compromised": false}'
            mock_client.chat.assert_called_once()

def test_get_inference_response_local_gguf_not_available():
    import backend.local_llm as real_llm
    with patch("backend.scanner.get_model_provider", return_value="local_gguf"):
        with patch.object(real_llm, "is_available", return_value=False):
            with pytest.raises(RuntimeError, match="not installed"):
                get_inference_response(
                    messages=[{"role": "user", "content": "test"}],
                    model_name="test"
                )


def test_get_inference_response_local_gguf_no_model():
    import backend.local_llm as real_llm
    with patch("backend.scanner.get_model_provider", return_value="local_gguf"):
        with patch.object(real_llm, "is_available", return_value=True):
            with patch.object(real_llm, "get_loaded_model_info", return_value=None):
                with pytest.raises(RuntimeError, match="no model is loaded"):
                    get_inference_response(
                        messages=[{"role": "user", "content": "test"}],
                        model_name="test"
                    )

def test_get_inference_response_local_gguf_success():
    import backend.local_llm as real_llm
    with patch("backend.scanner.get_model_provider", return_value="local_gguf"):
        with patch.object(real_llm, "is_available", return_value=True):
            with patch.object(real_llm, "get_loaded_model_info", return_value={"status": "loaded"}):
                with patch.object(real_llm, "chat_completion", return_value={"message": {"content": '{"compromised": true}'}}):
                    result = get_inference_response(
                        messages=[{"role": "user", "content": "test"}],
                        model_name="test",
                        options={"temperature": 0.0}
                    )
                    assert result["message"]["content"] == '{"compromised": true}'
