"""
Core PII scanning, pattern matching, OCR/multimodal inspection, and inference engine.

Provides:
- Multi-tier regex and entropy-based secret detection.
- Fast, cached filesystem traversal with incremental checksum change detection.
- Multimodal image and text inspection via Ollama or local GGUF models.
- Thread-safe scanner lifecycle management.
"""

import bisect
import concurrent.futures
import hashlib
import json
import math
import os
import re
import subprocess
import sys
import tempfile
import threading
import time
import urllib.request
from pathlib import Path
from typing import Any, Callable, Optional, Union

import ollama
from PIL import Image
import pillow_heif

try:
    from .state import load_cache, load_settings, save_cache
except (ImportError, ValueError):
    from backend.state import load_cache, load_settings, save_cache

# Register HEIF/HEIC decoder plugin for Pillow
pillow_heif.register_heif_opener()

# Default model identifiers — overridden by settings at runtime
DEFAULT_VISION_MODEL = "gemma4:12b"
DEFAULT_TEXT_MODEL = "gemma4:12b"

# File extension categories
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp"}
HEIC_EXTENSIONS = {".heic", ".heif"}
PDF_EXTENSIONS = {".pdf"}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".xls", ".pptx"}
TEXT_EXTENSIONS = {
    ".txt", ".csv", ".json", ".md", ".log", ".py",
    ".env", ".yaml", ".yml", ".ini", ".conf", ".properties", ".toml", ".sql",
    ".js", ".ts", ".jsx", ".tsx", ".sh", ".ps1", ".bat", ".html", ".htm", ".xml",
}

# Compiled regex patterns for PII and code secrets
PII_PATTERNS = {
    "SSN": re.compile(r"\b\d{3}-\d{2}-\d{4}\b"),
    "Credit Card": re.compile(r"\b(?:\d{4}[ -]?){3}\d{4}\b"),
}

VENDOR_PATTERNS = {
    "AWS Access Key": re.compile(r"\bAKIA[0-9A-Z]{16}\b"),
    "OpenAI API Key": re.compile(r"\bsk-[a-zA-Z0-9]{32,}\b|\bsk-proj-[a-zA-Z0-9_\-]{40,}\b"),
    "GitHub Token": re.compile(r"\bgh[po]_[a-zA-Z0-9]{36}\b"),
    "Private Key": re.compile(r"-----BEGIN (?:RSA|OPENSSH|EC) PRIVATE KEY-----"),
    "Database URI": re.compile(r"(?i)(?:mongodb|postgres|mysql|redis)://[^\s]+"),
}

GENERIC_SECRET_PATTERN = re.compile(r"\b[a-zA-Z0-9_\-\.\+=/]{16,}\b")
CONTEXT_PATTERN = re.compile(r"(?i)(secret|token|password|bearer|private)")

IMAGE_PROMPT = """
You are a data loss prevention assistant. Inspect this image carefully.
Does this image contain any sensitive or compromising Personally Identifiable Information (PII) such as:
- Social Security cards or numbers
- Credit cards, debit cards, or banking information
- Driver's licenses, passports, or government IDs
- Passwords, credentials, or private keys

Respond ONLY with JSON in this format:
{
  "compromised": true/false,
  "reason": "brief explanation of what was found",
  "items": [
    {
      "label": "Short label (e.g. SSN Card, Credit Card, Driver's License)",
      "box_2d": [ymin, xmin, ymax, xmax],
      "description": "Short description of detected item"
    }
  ]
}
Note: box_2d should use normalized coordinates from 0 to 1000 ([ymin, xmin, ymax, xmax]). If coordinates cannot be determined, set box_2d to null.
"""

TEXT_PROMPT = """
You are a data loss prevention assistant. Analyze the following text content.
Does it contain any compromising plaintext sensitive data such as Social Security Numbers (SSNs), credit card numbers, passwords, API secrets, or bank account details?

Respond ONLY with JSON in this format:
{
  "compromised": true/false,
  "reason": "brief explanation of what was found",
  "snippets": ["exact sensitive text snippet 1", "exact sensitive text snippet 2"]
}

Text Content:
"""

# Global synchronization lock for scanner state
state_lock = threading.Lock()


def get_file_type_label(ext: str) -> str:
    """Returns a standardized display label for a given file extension."""
    ext_lower = ext.lower()
    if ext_lower in IMAGE_EXTENSIONS:
        return "Image"
    elif ext_lower in HEIC_EXTENSIONS:
        return "HEIC"
    elif ext_lower in PDF_EXTENSIONS:
        return "PDF"
    elif ext_lower in OFFICE_EXTENSIONS:
        return "Office"
    return "Text"


def calculate_file_checksum(file_path: Union[str, Path], chunk_size: int = 65536) -> str:
    """
    Calculates SHA-256 checksum of a file in streaming chunks for memory efficiency.
    Returns hexadecimal digest string, or empty string on read/permission errors.
    """
    hasher = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            while chunk := f.read(chunk_size):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception as e:
        print(f"Error computing checksum for {file_path}: {e}")
        return ""


def get_active_vision_model() -> str:
    """Get the vision model name configured in settings."""
    settings = load_settings()
    return settings.get("vision_model_name", DEFAULT_VISION_MODEL)


def get_active_text_model() -> str:
    """Get the text model name configured in settings."""
    settings = load_settings()
    return settings.get("text_model_name", DEFAULT_TEXT_MODEL)


def get_model_provider() -> str:
    """Get the configured model provider: 'ollama' or 'local_gguf'."""
    settings = load_settings()
    return settings.get("model_provider", "ollama")


def get_inference_response(
    messages: list[dict[str, Any]],
    model_name: str,
    options: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    """
    Dispatches inference to the configured provider (Ollama or local GGUF).
    Returns a dict with {"message": {"content": "..."}} matching Ollama format.
    """
    provider = get_model_provider()

    if provider == "local_gguf":
        try:
            from . import local_llm
        except (ImportError, ValueError):
            import backend.local_llm as local_llm

        if not local_llm.is_available():
            raise RuntimeError(
                "Local GGUF mode is selected but llama-cpp-python is not installed. "
                "Install it with: pip install llama-cpp-python"
            )
        if local_llm.get_loaded_model_info() is None:
            raise RuntimeError(
                "Local GGUF mode is selected but no model is loaded. "
                "Go to Settings → Local Models and load a model first."
            )
        temperature = 0.0
        if options and "temperature" in options:
            temperature = options["temperature"]
        return local_llm.chat_completion(messages, temperature=temperature)
    else:
        client = get_client()
        kwargs: dict[str, Any] = {"model": model_name, "messages": messages}
        if options:
            kwargs["options"] = options
        return client.chat(**kwargs)


def calculate_shannon_entropy(data: str) -> float:
    """Calculates Shannon entropy for string randomness detection."""
    if not data:
        return 0.0
    entropy = 0.0
    for char in set(data):
        p_x = float(data.count(char)) / len(data)
        if p_x > 0:
            entropy -= p_x * math.log2(p_x)
    return entropy


def mask_secret(secret: str) -> str:
    """Masks high-entropy secrets for safe UI rendering."""
    if len(secret) <= 8:
        return "***"
    return secret[:8] + "..." + secret[-4:]


def get_line_starts(content: str) -> list[int]:
    """Computes byte/character offset start positions for every line."""
    return [0] + [m.end() for m in re.finditer(r"\n", content)]


def get_line_col(line_starts: list[int], index: int) -> tuple[int, int]:
    """Translates a flat string index into a 1-indexed (line_number, column_number)."""
    line_idx = bisect.bisect_right(line_starts, index) - 1
    if line_idx < 0:
        line_idx = 0
    return line_idx + 1, index - line_starts[line_idx]


def detect_secrets(content: str) -> list[dict[str, Any]]:
    """
    Detects API tokens, credentials, and high-entropy secrets in text content.
    Returns list of match dictionaries with line/col positions.
    """
    matches: list[dict[str, Any]] = []
    if not content:
        return matches

    line_starts = get_line_starts(content)
    tier1_intervals: list[tuple[int, int]] = []

    # Tier 1: Known vendor patterns (AWS, GitHub, OpenAI, DB URI, Private Keys)
    for pattern_name, pattern in VENDOR_PATTERNS.items():
        for m in pattern.finditer(content):
            start, end = m.start(), m.end()
            tier1_intervals.append((start, end))
            secret_val = m.group(0)
            masked_val = secret_val if pattern_name == "Private Key" else mask_secret(secret_val)

            line_num, start_col = get_line_col(line_starts, start)
            _, end_col = get_line_col(line_starts, end)

            matches.append({
                "line_number": line_num,
                "start_col": start_col,
                "end_col": end_col,
                "match_text": masked_val,
                "raw_match": secret_val,
                "pattern_name": pattern_name,
                "source": "regex",
            })

    # Tier 2 & 3: High-entropy tokens accompanied by contextual keywords
    for m in GENERIC_SECRET_PATTERN.finditer(content):
        start, end = m.start(), m.end()

        # Ignore if already matched by Tier 1
        if any(not (end <= t_start or start >= t_end) for t_start, t_end in tier1_intervals):
            continue

        candidate = m.group(0)
        if calculate_shannon_entropy(candidate) >= 4.5:
            context_start = max(0, start - 50)
            context_end = min(len(content), end + 50)
            context_window = content[context_start:context_end]

            if CONTEXT_PATTERN.search(context_window):
                line_num, start_col = get_line_col(line_starts, start)
                _, end_col = get_line_col(line_starts, end)
                matches.append({
                    "line_number": line_num,
                    "start_col": start_col,
                    "end_col": end_col,
                    "match_text": mask_secret(candidate),
                    "raw_match": candidate,
                    "pattern_name": "Generic API Token",
                    "source": "entropy_context",
                })

    return matches


def get_system_ram() -> int:
    """
    Get the total physical RAM of the system in bytes.
    Uses ctypes on Windows, sysctl on macOS, and /proc/meminfo on Linux.
    """
    try:
        if sys.platform == "win32":
            import ctypes

            class MEMORYSTATUSEX(ctypes.Structure):
                _fields_ = [
                    ("dwLength", ctypes.c_ulong),
                    ("dwMemoryLoad", ctypes.c_ulong),
                    ("ullTotalPhys", ctypes.c_ulonglong),
                    ("ullAvailPhys", ctypes.c_ulonglong),
                    ("ullTotalPageFile", ctypes.c_ulonglong),
                    ("ullAvailPageFile", ctypes.c_ulonglong),
                    ("ullTotalVirtual", ctypes.c_ulonglong),
                    ("ullAvailVirtual", ctypes.c_ulonglong),
                    ("ullAvailExtendedVirtual", ctypes.c_ulonglong),
                ]

            stat = MEMORYSTATUSEX()
            stat.dwLength = ctypes.sizeof(MEMORYSTATUSEX)
            ctypes.windll.kernel32.GlobalMemoryStatusEx(ctypes.byref(stat))
            return stat.ullTotalPhys
        elif sys.platform == "darwin":
            out = subprocess.check_output(["sysctl", "-n", "hw.memsize"])
            return int(out.strip())
        else:
            with open("/proc/meminfo", "r", encoding="utf-8") as f:
                for line in f:
                    if line.startswith("MemTotal:"):
                        return int(line.split()[1]) * 1024
    except Exception as e:
        print(f"Error checking RAM: {e}")
    return 8 * 1024 * 1024 * 1024  # Fallback to 8 GB


def get_auto_config(ram_bytes: int) -> dict[str, Any]:
    """Maps system RAM to recommended scanner concurrency and optimization settings."""
    gb = ram_bytes / (1024 ** 3)
    if gb < 8.0:
        return {"concurrency": 1, "image_optimization": "low", "text_scan_mode": "regex_llm"}
    elif gb < 16.0:
        return {"concurrency": 2, "image_optimization": "medium", "text_scan_mode": "regex_llm"}
    else:
        return {"concurrency": 4, "image_optimization": "medium", "text_scan_mode": "regex_llm"}


def get_optimized_image_path(image_path: str, optimization_setting: str) -> tuple[str, bool]:
    """
    Downsamples the image if needed, saving it to a temporary file.
    Returns (path_to_use, is_temporary).
    """
    if optimization_setting == "original":
        return image_path, False

    max_dim = 1024 if optimization_setting == "medium" else 768

    try:
        with Image.open(image_path) as img:
            w, h = img.size
            if w <= max_dim and h <= max_dim:
                return image_path, False

            suffix = Path(image_path).suffix.lower()
            if suffix not in {".jpg", ".jpeg", ".png", ".webp", ".bmp"}:
                suffix = ".jpg"

            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_img:
                temp_path = temp_img.name

            img.thumbnail((max_dim, max_dim), Image.Resampling.LANCZOS)
            if img.mode in ("RGBA", "P") and suffix in {".jpg", ".jpeg"}:
                img = img.convert("RGB")
            img.save(temp_path)
            return temp_path, True
    except Exception as e:
        print(f"Error optimizing image {image_path}: {e}")
        return image_path, False


def get_ollama_address() -> str:
    """Returns the configured Ollama server address stripped of trailing slashes."""
    settings = load_settings()
    return settings.get("ollama_address", "http://127.0.0.1:11434").rstrip("/")


def get_client() -> ollama.Client:
    """Returns an Ollama API client connected to the configured server."""
    return ollama.Client(host=get_ollama_address())


def ensure_ollama_running() -> tuple[bool, str]:
    """Checks if Ollama is accessible, attempting to spawn it locally if down."""
    addr = get_ollama_address()
    try:
        urllib.request.urlopen(addr, timeout=2)
        return True, "Ollama is running."
    except Exception:
        pass

    is_local = "127.0.0.1" in addr or "localhost" in addr
    if not is_local:
        return False, f"Could not connect to remote Ollama server at {addr}."

    try:
        popen_kwargs: dict[str, Any] = {}
        if os.name == "nt" and hasattr(subprocess, "CREATE_NO_WINDOW"):
            popen_kwargs["creationflags"] = subprocess.CREATE_NO_WINDOW

        settings = load_settings()
        concurrency_setting = settings.get("concurrency", "auto")
        if concurrency_setting == "auto":
            ram = get_system_ram()
            concurrency = get_auto_config(ram)["concurrency"]
        else:
            try:
                concurrency = int(concurrency_setting)
            except ValueError:
                concurrency = 2

        env = os.environ.copy()
        env["OLLAMA_NUM_PARALLEL"] = str(concurrency)
        env["OLLAMA_MAX_LOADED_MODELS"] = str(max(concurrency, 2))
        popen_kwargs["env"] = env

        subprocess.Popen(["ollama", "serve"], **popen_kwargs)

        for _ in range(10):
            time.sleep(1)
            try:
                urllib.request.urlopen(addr, timeout=2)
                return True, "Ollama started successfully."
            except Exception:
                pass

        return False, "Failed to start Ollama. The process launched but the endpoint is unresponsive."
    except FileNotFoundError:
        return False, "Ollama executable not found. Please ensure Ollama is installed and in your PATH."
    except Exception as e:
        return False, f"Unexpected error starting Ollama: {str(e)}"


# ---------------------------------------------------------------------------
# Document Text Extraction Helpers
# ---------------------------------------------------------------------------

def _extract_pdf_text(path: Path) -> str:
    import pypdf
    reader = pypdf.PdfReader(path)
    parts = [page.extract_text() for page in reader.pages[:15] if page.extract_text()]
    return "\n".join(parts)


def _extract_docx_text(path: Path) -> str:
    import docx
    doc = docx.Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    try:
        for sheet in wb.worksheets[:5]:
            for row in sheet.iter_rows(values_only=True, max_row=100):
                row_vals = [str(val).strip() for val in row if val is not None and str(val).strip()]
                if row_vals:
                    parts.append(" | ".join(row_vals))
    finally:
        wb.close()
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    import pptx
    prs = pptx.Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def get_file_text_content(file_path: Union[str, Path]) -> str:
    """Extracts up to 4096 characters of text from PDF, DOCX, XLSX, PPTX, or text files."""
    path = Path(file_path)
    ext = path.suffix.lower()
    try:
        if ext in PDF_EXTENSIONS:
            text = _extract_pdf_text(path)
        elif ext == ".docx":
            text = _extract_docx_text(path)
        elif ext in (".xlsx", ".xls"):
            text = _extract_xlsx_text(path)
        elif ext == ".pptx":
            text = _extract_pptx_text(path)
        else:
            text = path.read_text(errors="ignore")
        return text[:4096]
    except Exception as e:
        print(f"Error extracting text from {path}: {e}")
        return ""


# ---------------------------------------------------------------------------
# Scanner State & AI Output Parsing
# ---------------------------------------------------------------------------

class ScannerState:
    """Thread-safe telemetry and progress state container."""

    def __init__(self) -> None:
        self.is_scanning: bool = False
        self.progress: dict[str, Any] = {
            "total_files": 0,
            "scanned_files": 0,
            "current_file": "",
            "flagged_count": 0,
            "skipped_count": 0,
        }
        self.flagged_files: list[dict[str, Any]] = []
        self.should_stop: bool = False

    def reset(self) -> None:
        self.is_scanning = False
        self.progress = {
            "total_files": 0,
            "scanned_files": 0,
            "current_file": "",
            "flagged_count": 0,
            "skipped_count": 0,
        }
        self.flagged_files = []
        self.should_stop = False


scan_state = ScannerState()


def parse_ai_response(raw_output: str) -> dict[str, Any]:
    """
    Parses LLM JSON response, extracting compromised status, reason, items, and snippets.
    Provides robust fallback regex parsing for markdown-wrapped or malformed outputs.
    """
    try:
        cleaned = raw_output.strip().removeprefix("```json").removesuffix("```").strip()
        data = json.loads(cleaned)
        if not isinstance(data, dict):
            raise ValueError("Expected JSON object")

        compromised = bool(data.get("compromised", False))
        reason = str(data.get("reason", "Flagged by AI analysis"))
        items = data.get("items", []) if isinstance(data.get("items"), list) else []
        snippets = data.get("snippets", []) if isinstance(data.get("snippets"), list) else []

        valid_items = []
        for item in items:
            if isinstance(item, dict):
                label = str(item.get("label", "Sensitive Item"))
                box = item.get("box_2d")
                if isinstance(box, list) and len(box) == 4 and all(isinstance(x, (int, float)) for x in box):
                    valid_box = [max(0, min(1000, float(x))) for x in box]
                else:
                    valid_box = None
                valid_items.append({
                    "label": label,
                    "box_2d": valid_box,
                    "description": str(item.get("description", label)),
                })

        return {
            "compromised": compromised,
            "reason": reason,
            "items": valid_items,
            "snippets": [str(s) for s in snippets if s],
        }
    except Exception:
        # Regex JSON fallback
        json_match = re.search(r"\{[\s\S]*\}", raw_output)
        if json_match:
            try:
                data = json.loads(json_match.group(0))
                if isinstance(data, dict):
                    return {
                        "compromised": bool(data.get("compromised", False)),
                        "reason": str(data.get("reason", "Flagged by AI analysis")),
                        "items": data.get("items", []) if isinstance(data.get("items"), list) else [],
                        "snippets": data.get("snippets", []) if isinstance(data.get("snippets"), list) else [],
                    }
            except Exception:
                pass

        is_flagged = "true" in raw_output.lower() and "compromised" in raw_output.lower()
        return {
            "compromised": is_flagged,
            "reason": raw_output.replace("\n", " ").strip(),
            "items": [],
            "snippets": [],
        }


def locate_text_pii_matches(
    content: str,
    ai_snippets: Optional[list[str]] = None,
    file_path: Optional[str] = None,
) -> list[dict[str, Any]]:
    """
    Locates exact line numbers, character positions, and match strings
    for regex patterns and AI-flagged snippets in text content.
    Returns a list of match dicts sorted by line number and column.
    """
    if not content:
        return []

    matches: list[dict[str, Any]] = []

    try:
        matches.extend(detect_secrets(content))
    except Exception as e:
        print(f"Error in detect_secrets: {e}")

    lines = content.splitlines()

    # 1. Regex pattern matches (SSN, Credit Cards)
    for line_idx, line in enumerate(lines):
        line_num = line_idx + 1
        for pattern_name, pattern in PII_PATTERNS.items():
            for m in pattern.finditer(line):
                matches.append({
                    "line_number": line_num,
                    "start_col": m.start(),
                    "end_col": m.end(),
                    "match_text": m.group(0),
                    "pattern_name": pattern_name,
                    "source": "regex",
                })

    # 2. AI snippets matches
    if ai_snippets:
        for snippet in ai_snippets:
            if not isinstance(snippet, str) or not snippet.strip():
                continue
            snippet_clean = snippet.strip()
            lower_snippet = snippet_clean.lower()
            for line_idx, line in enumerate(lines):
                line_num = line_idx + 1
                lower_line = line.lower()
                start_pos = 0
                while True:
                    pos = lower_line.find(lower_snippet, start_pos)
                    if pos == -1:
                        break

                    already_covered = any(
                        m["line_number"] == line_num
                        and (pos <= m["start_col"] < pos + len(snippet_clean) or m["start_col"] <= pos < m["end_col"])
                        for m in matches
                    )
                    if not already_covered:
                        matches.append({
                            "line_number": line_num,
                            "start_col": pos,
                            "end_col": pos + len(snippet_clean),
                            "match_text": line[pos:pos + len(snippet_clean)],
                            "pattern_name": "AI Finding",
                            "source": "ai",
                        })
                    start_pos = pos + len(lower_snippet)
                    if start_pos >= len(lower_line):
                        break

    # 3. Filter whitelisted exceptions
    if file_path:
        try:
            try:
                from .remediation import is_file_or_match_ignored
            except (ImportError, ValueError):
                from backend.remediation import is_file_or_match_ignored
            matches = [
                m for m in matches
                if not (
                    is_file_or_match_ignored(file_path, match_text=m.get("match_text"), pattern_name=m.get("pattern_name"))
                    or (m.get("raw_match") and is_file_or_match_ignored(file_path, match_text=m.get("raw_match"), pattern_name=m.get("pattern_name")))
                )
            ]
        except Exception:
            pass

    matches.sort(key=lambda m: (m["line_number"], m["start_col"]))
    return matches


# ---------------------------------------------------------------------------
# Inspection & Inference
# ---------------------------------------------------------------------------

def inspect_image(image_path: str, optimization_setting: str = "medium") -> dict[str, Any]:
    """Inspects an image file for sensitive items using vision models."""
    temp_path = None
    is_temp = False
    try:
        path_to_use, is_temp = get_optimized_image_path(image_path, optimization_setting)
        if is_temp:
            temp_path = path_to_use

        provider = get_model_provider()
        vision_model = get_active_vision_model()

        if provider == "local_gguf":
            try:
                from . import local_llm
            except (ImportError, ValueError):
                import backend.local_llm as local_llm
            if not local_llm.is_available() or local_llm.get_loaded_model_info() is None:
                return {"compromised": False, "reason": "Local model not loaded for image inspection"}
            response = local_llm.chat_completion(
                [{"role": "user", "content": IMAGE_PROMPT + "\n[Image provided at: " + path_to_use + "]"}],
                temperature=0.0,
            )
        else:
            client = get_client()
            response = client.chat(
                model=vision_model,
                messages=[
                    {
                        "role": "user",
                        "content": IMAGE_PROMPT,
                        "images": [path_to_use],
                    }
                ],
                options={"temperature": 0.0},
            )
        return parse_ai_response(response["message"]["content"])
    except Exception as e:
        return {"compromised": False, "reason": f"Error inspecting image: {str(e)}"}
    finally:
        if is_temp and temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception as e:
                print(f"Error removing optimized image temp file {temp_path}: {e}")


def process_heic_image(file_path: Union[str, Path], optimization_setting: str = "medium") -> dict[str, Any]:
    """Converts a HEIC/HEIF image to JPEG and inspects it."""
    temp_path = None
    try:
        with Image.open(file_path) as img:
            max_dim = 1024 if optimization_setting == "medium" else 768 if optimization_setting == "low" else None
            if max_dim:
                img.thumbnail((max_dim, max_dim), Image.Resampling.LANCZOS)

            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as temp_img:
                temp_path = temp_img.name
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                img.save(temp_path, "JPEG")

            return inspect_image(temp_path, optimization_setting="original")
    except Exception as e:
        return {"compromised": False, "reason": f"Error processing HEIC file: {str(e)}"}
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception as e:
                print(f"Error removing HEIC temp file {temp_path}: {e}")


def inspect_text(file_path: Union[str, Path], text_scan_mode: str = "regex_llm") -> dict[str, Any]:
    """Inspects text-based files using regex filtering and/or AI verification."""
    try:
        content = get_file_text_content(file_path)
        if not content.strip():
            return {"compromised": False, "reason": "Empty or unreadable file"}

        matched_patterns: list[str] = []
        if text_scan_mode in ("regex_llm", "regex_only"):
            for name, pattern in PII_PATTERNS.items():
                if pattern.search(content):
                    matched_patterns.append(name)

            try:
                for sm in detect_secrets(content):
                    if sm["pattern_name"] not in matched_patterns:
                        matched_patterns.append(sm["pattern_name"])
            except Exception:
                pass

            if not matched_patterns:
                return {"compromised": False, "reason": "No PII patterns matched via regex"}

            if text_scan_mode == "regex_only":
                return {
                    "compromised": True,
                    "reason": f"Flagged via regex filter: matched {', '.join(matched_patterns)}",
                }

            if text_scan_mode == "regex_llm":
                return {
                    "compromised": True,
                    "reason": f"Flagged via regex filter: matched {', '.join(matched_patterns)}. Awaiting AI Verification.",
                    "needs_ai_verification": True,
                }

        text_model = get_active_text_model()
        response = get_inference_response(
            messages=[{"role": "user", "content": TEXT_PROMPT + "\n" + content}],
            model_name=text_model,
            options={"temperature": 0.0},
        )
        return parse_ai_response(response["message"]["content"])
    except Exception as e:
        return {"compromised": False, "reason": f"Error inspecting text: {str(e)}"}


def verify_text_file_with_ai(file_path: Union[str, Path]) -> dict[str, Any]:
    """Explicitly triggers full AI verification on a regex-flagged text file."""
    try:
        path = Path(file_path)
        content = get_file_text_content(path)
        if not content.strip():
            return {"compromised": False, "reason": "Empty file"}

        text_model = get_active_text_model()
        response = get_inference_response(
            messages=[{"role": "user", "content": TEXT_PROMPT + "\\n" + content}],
            model_name=text_model,
            options={"temperature": 0.0},
        )
        return parse_ai_response(response["message"]["content"])
    except Exception as e:
        return {"compromised": False, "reason": f"Error verifying text: {str(e)}"}


# ---------------------------------------------------------------------------
# Directory Traversal & Scan Execution
# ---------------------------------------------------------------------------

def get_scannable_files(folders: list[str]) -> list[Path]:
    """Discovers all supported image, document, and text files within target folders."""
    files_to_scan: list[Path] = []
    try:
        from .remediation import is_file_or_match_ignored
    except (ImportError, ValueError):
        try:
            from backend.remediation import is_file_or_match_ignored
        except Exception:
            is_file_or_match_ignored = lambda f: False

    supported_extensions = IMAGE_EXTENSIONS | HEIC_EXTENSIONS | PDF_EXTENSIONS | OFFICE_EXTENSIONS | TEXT_EXTENSIONS

    for folder in folders:
        path = Path(folder)
        if not path.exists() or not path.is_dir():
            continue

        for root, _, files in os.walk(path):
            for file in files:
                file_path = Path(root) / file
                if file_path.suffix.lower() in supported_extensions:
                    if not is_file_or_match_ignored(str(file_path)):
                        files_to_scan.append(file_path)
    return files_to_scan


def run_scan(
    folders: list[str],
    save_results_callback: Callable[[list[dict[str, Any]]], None],
    rescan_all: bool = False,
    scan_id: Optional[int] = None,
) -> None:
    """
    Executes a parallel filesystem scan across target directories.
    Handles incremental checksum caching, thread limits, and real-time state updates.
    """
    global _scan_id
    if scan_id is None:
        scan_id = _scan_id

    with state_lock:
        scan_state.reset()
        scan_state.is_scanning = True
        scan_state.should_stop = False

    file_cache = load_cache()
    settings = load_settings()

    concurrency_setting = settings.get("concurrency", "auto")
    image_opt = settings.get("image_optimization", "medium")
    text_mode = settings.get("text_scan_mode", "regex_llm")
    auto_delete = settings.get("auto_delete", False)

    if concurrency_setting == "auto":
        concurrency = get_auto_config(get_system_ram())["concurrency"]
    else:
        try:
            concurrency = int(concurrency_setting)
        except ValueError:
            concurrency = 2

    files = get_scannable_files(folders)
    scan_state.progress["total_files"] = len(files)

    def scan_single_file(file_path: Path) -> None:
        if scan_state.should_stop or scan_id != _scan_id:
            return

        with state_lock:
            scan_state.progress["current_file"] = str(file_path)

        str_path = str(file_path)
        try:
            stat = os.stat(str_path)
            mtime = stat.st_mtime
            size = stat.st_size
        except OSError:
            mtime = 0
            size = 0

        current_checksum = calculate_file_checksum(str_path)
        ext = file_path.suffix.lower()
        file_type_label = get_file_type_label(ext)

        # Check incremental cache
        if not rescan_all and str_path in file_cache:
            cached_entry = file_cache[str_path]
            cached_checksum = cached_entry.get("checksum")

            is_unmodified = False
            if cached_checksum:
                is_unmodified = bool(current_checksum and cached_checksum == current_checksum)
            elif cached_entry.get("mtime") == mtime:
                is_unmodified = True
                cached_entry["checksum"] = current_checksum
                cached_entry["size"] = size

            if is_unmodified:
                cached_result = cached_entry.get("result")
                if not cached_result or not cached_result.get("compromised"):
                    with state_lock:
                        scan_state.progress["skipped_count"] += 1
                        scan_state.progress["scanned_files"] += 1
                    return
                else:
                    with state_lock:
                        scan_state.flagged_files.append({
                            "file": str_path,
                            "type": file_type_label,
                            "reason": cached_result.get("reason"),
                            "selected": False,
                            "auto_deleted": False,
                            "needs_ai_verification": cached_result.get("needs_ai_verification", False),
                            "compromised": True,
                            "items": cached_result.get("items", []),
                            "snippets": cached_result.get("snippets", []),
                            "checksum": current_checksum,
                        })
                        scan_state.progress["flagged_count"] = len(scan_state.flagged_files)
                        scan_state.progress["scanned_files"] += 1
                    return

        # Perform file inspection
        result = None
        try:
            if ext in IMAGE_EXTENSIONS:
                result = inspect_image(str(file_path), image_opt)
            elif ext in HEIC_EXTENSIONS:
                result = process_heic_image(file_path, image_opt)
            elif ext in PDF_EXTENSIONS or ext in OFFICE_EXTENSIONS or ext in TEXT_EXTENSIONS:
                result = inspect_text(file_path, text_mode)
        except Exception as e:
            print(f"Error scanning file {file_path}: {e}")

        if scan_state.should_stop or scan_id != _scan_id:
            return

        file_cache[str_path] = {
            "mtime": mtime,
            "size": size,
            "checksum": current_checksum,
            "result": result,
        }

        if result and result.get("compromised"):
            deleted_successfully = False
            needs_ai_verification = result.get("needs_ai_verification", False)
            if auto_delete and not needs_ai_verification:
                try:
                    if os.path.exists(file_path):
                        os.remove(file_path)
                        deleted_successfully = True
                except Exception as e:
                    print(f"Failed to auto-delete {file_path}: {e}")

            with state_lock:
                scan_state.flagged_files.append({
                    "file": str(file_path),
                    "type": file_type_label,
                    "reason": result.get("reason"),
                    "selected": False,
                    "auto_deleted": deleted_successfully,
                    "needs_ai_verification": needs_ai_verification,
                    "compromised": True,
                    "items": result.get("items", []),
                    "snippets": result.get("snippets", []),
                    "checksum": current_checksum,
                })
                scan_state.progress["flagged_count"] = len(scan_state.flagged_files)

        with state_lock:
            scan_state.progress["scanned_files"] += 1

    executor = concurrent.futures.ThreadPoolExecutor(max_workers=concurrency)
    try:
        futures = {executor.submit(scan_single_file, f): f for f in files}
        for future in concurrent.futures.as_completed(futures):
            if scan_state.should_stop or scan_id != _scan_id:
                try:
                    executor.shutdown(wait=False, cancel_futures=True)
                except TypeError:
                    executor.shutdown(wait=False)
                break

        if not scan_state.should_stop and scan_id == _scan_id:
            save_cache(file_cache)
            save_results_callback(scan_state.flagged_files)
    finally:
        try:
            executor.shutdown(wait=False, cancel_futures=True)
        except Exception:
            pass
        if scan_id == _scan_id:
            with state_lock:
                scan_state.is_scanning = False


_scan_id: int = 0
_scan_lock = threading.Lock()
_current_scan_thread: Optional[threading.Thread] = None


def start_scan_thread(
    folders: list[str],
    save_results_callback: Callable[[list[dict[str, Any]]], None],
    rescan_all: bool = False,
) -> bool:
    """Spawns a background scanning thread, preventing duplicate concurrent runs."""
    global _scan_id, _current_scan_thread
    with _scan_lock:
        if scan_state.is_scanning and not scan_state.should_stop:
            return False

        _scan_id += 1
        current_id = _scan_id

        with state_lock:
            scan_state.reset()
            scan_state.is_scanning = True
            scan_state.should_stop = False

        _current_scan_thread = threading.Thread(
            target=run_scan,
            args=(folders, save_results_callback, rescan_all, current_id),
            name=f"ArgusScanWorker-{current_id}",
            daemon=True,
        )
        _current_scan_thread.start()
        return True


def stop_scan(timeout: Optional[float] = None) -> None:
    """Signals the active scan thread to terminate gracefully."""
    global _current_scan_thread
    with state_lock:
        scan_state.should_stop = True
        scan_state.is_scanning = False
    if timeout and _current_scan_thread and _current_scan_thread.is_alive():
        _current_scan_thread.join(timeout=timeout)
