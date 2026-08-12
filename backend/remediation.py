import os
import sys
import stat
import json
import time
import shutil
import uuid
import re
from datetime import datetime
from pathlib import Path

try:
    from .state import load_settings, save_settings, load_results, save_results, load_cache, save_cache
    from .scanner import calculate_file_checksum, locate_text_pii_matches, scan_state, state_lock
except (ImportError, ValueError):
    from backend.state import load_settings, save_settings, load_results, save_results, load_cache, save_cache
    from backend.scanner import calculate_file_checksum, locate_text_pii_matches, scan_state, state_lock

BACKUP_DIR_NAME = ".argus_backups"
ARGUSIGNORE_FILENAME = ".argusignore"
TRASH_DIR_NAME = ".argus_trash"


# --------------------------------------------------------------------------
# Permissions & Integrity Checks
# --------------------------------------------------------------------------

def check_write_permission(file_path: str) -> bool:
    """Check if the given file exists and has write permission."""
    if not os.path.exists(file_path):
        return False
    if not os.access(file_path, os.W_OK):
        return False
    try:
        st = os.stat(file_path)
        if sys.platform == "win32":
            return bool(st.st_mode & stat.S_IWRITE)
        return True
    except Exception:
        return False


def fix_file_permissions(file_path: str) -> dict:
    """
    Remove read-only attributes from a file to enable in-place remediation.
    On Windows, uses ctypes/os.chmod to clear the readonly flag.
    """
    if not os.path.exists(file_path):
        return {"success": False, "error": "File not found", "file": file_path}
    try:
        # Clear read-only flags
        os.chmod(file_path, stat.S_IWRITE | stat.S_IREAD)
        if sys.platform == "win32":
            try:
                import ctypes
                FILE_ATTRIBUTE_NORMAL = 0x80
                ctypes.windll.kernel32.SetFileAttributesW(str(file_path), FILE_ATTRIBUTE_NORMAL)
            except Exception:
                pass
        is_writable = check_write_permission(file_path)
        return {"success": is_writable, "file": file_path, "writable": is_writable}
    except Exception as e:
        return {"success": False, "error": str(e), "file": file_path}


def verify_file_integrity(file_path: str, expected_checksum: str = None) -> tuple[bool, str]:
    """
    Checks if a file exists and whether its current checksum matches expected_checksum.
    Returns (is_valid, current_checksum_or_error).
    """
    if not os.path.exists(file_path):
        return False, "file_not_found"
    current_checksum = calculate_file_checksum(file_path)
    if expected_checksum and expected_checksum != current_checksum:
        return False, "checksum_mismatch"
    return True, current_checksum


# --------------------------------------------------------------------------
# Backups Management (.argus_backups/)
# --------------------------------------------------------------------------

def _get_backup_dir(base_dir: str = ".") -> str:
    """Resolve and ensure the .argus_backups directory exists."""
    backup_dir = os.path.join(base_dir, BACKUP_DIR_NAME)
    os.makedirs(backup_dir, exist_ok=True)
    return backup_dir


def _get_backup_index_path(base_dir: str = ".") -> str:
    return os.path.join(_get_backup_dir(base_dir), "index.json")


def _load_backup_index(base_dir: str = ".") -> list[dict]:
    index_path = _get_backup_index_path(base_dir)
    if os.path.exists(index_path):
        try:
            with open(index_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return []


def _save_backup_index(index_data: list[dict], base_dir: str = "."):
    index_path = _get_backup_index_path(base_dir)
    with open(index_path, "w", encoding="utf-8") as f:
        json.dump(index_data, f, indent=2)


def create_backup(file_path: str, reason: str = "Pre-redaction backup", base_dir: str = ".") -> dict:
    """
    Copies the original file to .argus_backups/ before modifying it.
    Maintains a metadata index and automatically runs 7-day retention pruning.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File to backup not found: {file_path}")

    backup_dir = _get_backup_dir(base_dir)
    prune_expired_backups(max_days=7, base_dir=base_dir)

    timestamp = time.time()
    date_str = datetime.fromtimestamp(timestamp).strftime("%Y%m%d_%H%M%S")
    orig_name = os.path.basename(file_path)
    backup_id = str(uuid.uuid4())[:8]
    backup_filename = f"{date_str}_{backup_id}_{orig_name}"
    backup_full_path = os.path.join(backup_dir, backup_filename)

    shutil.copy2(file_path, backup_full_path)
    checksum = calculate_file_checksum(file_path)
    size = os.path.getsize(file_path)

    entry = {
        "id": backup_id,
        "original_path": os.path.abspath(file_path),
        "backup_path": os.path.abspath(backup_full_path),
        "filename": orig_name,
        "backup_filename": backup_filename,
        "timestamp": timestamp,
        "created_at": datetime.fromtimestamp(timestamp).strftime("%Y-%m-%d %H:%M:%S"),
        "checksum": checksum,
        "size": size,
        "reason": reason
    }

    index = _load_backup_index(base_dir)
    index.insert(0, entry)
    _save_backup_index(index, base_dir)

    return entry


def prune_expired_backups(max_days: int = 7, base_dir: str = ".") -> int:
    """
    Removes backups older than max_days (default 7 days) from disk and index.
    Returns the count of pruned backup files.
    """
    backup_dir = _get_backup_dir(base_dir)
    index = _load_backup_index(base_dir)
    now = time.time()
    cutoff_time = now - (max_days * 86400)

    kept_entries = []
    pruned_count = 0

    for entry in index:
        entry_time = entry.get("timestamp", 0)
        backup_path = entry.get("backup_path", "")
        if entry_time < cutoff_time:
            if os.path.exists(backup_path):
                try:
                    os.remove(backup_path)
                except Exception as e:
                    print(f"Error removing expired backup {backup_path}: {e}")
            pruned_count += 1
        else:
            kept_entries.append(entry)

    _save_backup_index(kept_entries, base_dir)
    return pruned_count


def list_backups(base_dir: str = ".") -> list[dict]:
    """Returns list of backups sorted newest to oldest."""
    return _load_backup_index(base_dir)


def restore_backup(backup_id_or_path: str, base_dir: str = ".") -> dict:
    """
    Restores the original file from a backup file.
    """
    index = _load_backup_index(base_dir)
    target_entry = None
    for entry in index:
        if entry.get("id") == backup_id_or_path or entry.get("backup_path") == backup_id_or_path:
            target_entry = entry
            break

    if not target_entry:
        return {"success": False, "error": "Backup record not found"}

    backup_path = target_entry["backup_path"]
    original_path = target_entry["original_path"]

    if not os.path.exists(backup_path):
        return {"success": False, "error": "Backup file missing on disk"}

    try:
        shutil.copy2(backup_path, original_path)
        new_checksum = calculate_file_checksum(original_path)
        
        # Clear cache entry to force re-evaluation
        cache = load_cache()
        if original_path in cache:
            del cache[original_path]
            save_cache(cache)

        return {
            "success": True,
            "restored_file": original_path,
            "new_checksum": new_checksum
        }
    except Exception as e:
        return {"success": False, "error": f"Failed to restore: {str(e)}"}


# --------------------------------------------------------------------------
# Masking & Redaction Logic
# --------------------------------------------------------------------------

def mask_text(match_text: str, mask_pattern: str = "redacted") -> str:
    """
    Formats replacement text based on pattern:
    - 'redacted': '[REDACTED]'
    - 'mask' / 'asterisk': maintains format structure (e.g. 123-45-6789 -> XXX-XX-6789)
    - 'confidential': '[CONFIDENTIAL]'
    - custom: returns custom mask
    """
    if not match_text:
        return "[REDACTED]"

    if mask_pattern == "redacted":
        return "[REDACTED]"
    elif mask_pattern == "confidential":
        return "[CONFIDENTIAL]"
    elif mask_pattern in ("mask", "asterisk"):
        # Detect SSN pattern
        if re.match(r"^\d{3}-\d{2}-\d{4}$", match_text):
            return f"XXX-XX-{match_text[-4:]}"
        # Detect Credit card pattern
        if re.match(r"^(?:\d{4}[ -]?){3}\d{4}$", match_text):
            digits = re.sub(r"\D", "", match_text)
            return f"XXXX-XXXX-XXXX-{digits[-4:]}"
        # Generic alphanumeric masking: replace alphanumeric with '*'
        return re.sub(r"[A-Za-z0-9]", "*", match_text)
    elif mask_pattern:
        return str(mask_pattern)
    return "[REDACTED]"


def _redact_plain_text(content: str, line_number: int, start_col: int, end_col: int, match_text: str, masked: str) -> str:
    """Replaces match_text on the specified line with masked text."""
    lines = content.splitlines(keepends=True)
    if not lines:
        return content

    target_idx = line_number - 1
    if 0 <= target_idx < len(lines):
        line = lines[target_idx]
        
        # Check if line substring matches at exact start_col/end_col
        if 0 <= start_col < len(line) and line[start_col:end_col] == match_text:
            new_line = line[:start_col] + masked + line[end_col:]
            lines[target_idx] = new_line
        else:
            # Handle masked match_text (e.g. from secret detection)
            if "..." in match_text and 0 <= start_col < end_col <= len(line):
                seg = line[start_col:end_col]
                # Allow redaction if the original line segment matches the unmasked parts
                if match_text.startswith(seg[:4]) or match_text.endswith(seg[-4:]):
                    new_line = line[:start_col] + masked + line[end_col:]
                    lines[target_idx] = new_line
                    return "".join(lines)
                    
            if match_text in line:
                # Fallback: replace first occurrence of match_text on this line
                lines[target_idx] = line.replace(match_text, masked, 1)
            else:
                # Global single replacement fallback
                full = "".join(lines)
                if match_text in full:
                    return full.replace(match_text, masked, 1)
        return "".join(lines)
    else:
        # Fallback if line index is out of range
        full = "".join(lines)
        return full.replace(match_text, masked, 1)


def redact_file_entity(
    file_path: str,
    line_number: int = 1,
    start_col: int = 0,
    end_col: int = 0,
    match_text: str = "",
    mask_pattern: str = "redacted",
    expected_checksum: str = None,
    base_dir: str = "."
) -> dict:
    """
    Sanitizes a file in-place by masking only the specific detected PII entity.
    Generates an automatic backup in .argus_backups/ and validates checksum against external edits.
    """
    if not os.path.exists(file_path):
        return {"success": False, "error": "file_not_found", "message": f"File not found: {file_path}"}

    if not check_write_permission(file_path):
        return {
            "success": False,
            "error": "permission_denied",
            "message": "Permission Denied: File is read-only or lacks write permission.",
            "file": file_path
        }

    is_valid, current_checksum = verify_file_integrity(file_path, expected_checksum)
    if not is_valid:
        return {
            "success": False,
            "error": "file_modified",
            "message": "File was modified externally on disk. Preview reloaded with fresh content.",
            "file": file_path,
            "current_checksum": current_checksum
        }

    ext = os.path.splitext(file_path)[1].lower()
    from .scanner import IMAGE_EXTENSIONS, HEIC_EXTENSIONS, PDF_EXTENSIONS, OFFICE_EXTENSIONS

    if ext in IMAGE_EXTENSIONS or ext in HEIC_EXTENSIONS:
        return {
            "success": False,
            "error": "unsupported_format",
            "message": "In-place text redaction is not supported for image formats. Please use Delete File or Mark as Safe.",
            "file": file_path
        }

    masked = mask_text(match_text, mask_pattern)
    backup_entry = create_backup(file_path, reason=f"Redact: {match_text[:20]}", base_dir=base_dir)

    try:
        if ext == ".docx":
            import docx
            doc = docx.Document(file_path)
            for p in doc.paragraphs:
                if match_text in p.text:
                    p.text = p.text.replace(match_text, masked)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if match_text in cell.text:
                            cell.text = cell.text.replace(match_text, masked)
            doc.save(file_path)

        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str) and match_text in cell.value:
                            cell.value = cell.value.replace(match_text, masked)
            wb.save(file_path)
            wb.close()

        elif ext == ".pptx":
            import pptx
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text and match_text in shape.text:
                        for p in shape.text_frame.paragraphs:
                            if match_text in p.text:
                                p.text = p.text.replace(match_text, masked)
            prs.save(file_path)

        elif ext in PDF_EXTENSIONS:
            # Direct in-place PDF text mutation is not supported without reflowing
            return {
                "success": False,
                "error": "unsupported_format",
                "message": "Direct in-place text replacement is not supported for PDF documents. Please use Delete File or Mark as Safe.",
                "file": file_path
            }

        else:
            # Plain text / code / config files
            encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
            content = None
            used_encoding = "utf-8"
            for enc in encodings:
                try:
                    with open(file_path, "r", encoding=enc) as f:
                        content = f.read()
                        used_encoding = enc
                        break
                except UnicodeDecodeError:
                    continue

            if content is None:
                with open(file_path, "r", errors="ignore") as f:
                    content = f.read()

            updated_content = _redact_plain_text(content, line_number, start_col, end_col, match_text, masked)
            with open(file_path, "w", encoding=used_encoding) as f:
                f.write(updated_content)

        new_checksum = calculate_file_checksum(file_path)

        # Update cache & results
        _sync_state_after_remediation(file_path, new_checksum)

        return {
            "success": True,
            "file": file_path,
            "new_checksum": new_checksum,
            "backup_path": backup_entry["backup_path"],
            "masked": masked
        }
    except Exception as e:
        return {"success": False, "error": str(e), "message": f"Error during redaction: {str(e)}"}


def batch_redact_file(
    file_path: str,
    highlights: list[dict] = None,
    mask_pattern: str = "redacted",
    expected_checksum: str = None,
    base_dir: str = "."
) -> dict:
    """
    Redacts all provided highlights in a single file in a single atomic pass.
    Sorts highlights from bottom-to-top, right-to-left to preserve character column offsets.
    """
    if not os.path.exists(file_path):
        return {"success": False, "error": "file_not_found", "message": f"File not found: {file_path}"}

    if not check_write_permission(file_path):
        return {
            "success": False,
            "error": "permission_denied",
            "message": "Permission Denied: File is read-only.",
            "file": file_path
        }

    is_valid, current_checksum = verify_file_integrity(file_path, expected_checksum)
    if not is_valid:
        return {
            "success": False,
            "error": "file_modified",
            "message": "File was modified externally on disk.",
            "file": file_path,
            "current_checksum": current_checksum
        }

    ext = os.path.splitext(file_path)[1].lower()
    from .scanner import IMAGE_EXTENSIONS, HEIC_EXTENSIONS, PDF_EXTENSIONS

    if ext in IMAGE_EXTENSIONS or ext in HEIC_EXTENSIONS or ext in PDF_EXTENSIONS:
        return {
            "success": False,
            "error": "unsupported_format",
            "message": f"Batch in-place text redaction is not supported for {ext} files.",
            "file": file_path
        }

    # If highlights not given, locate them automatically
    if not highlights:
        from .scanner import get_file_text_content
        text = get_file_text_content(Path(file_path))
        highlights = locate_text_pii_matches(text)

    if not highlights:
        return {"success": True, "file": file_path, "redacted_count": 0, "message": "No findings to redact."}

    backup_entry = create_backup(file_path, reason=f"Batch redact {len(highlights)} entities", base_dir=base_dir)

    try:
        # Read text content
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
        content = None
        used_encoding = "utf-8"
        for enc in encodings:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    content = f.read()
                    used_encoding = enc
                    break
            except UnicodeDecodeError:
                continue
        if content is None:
            with open(file_path, "r", errors="ignore") as f:
                content = f.read()

        # Sort reverse by line number and column
        sorted_highlights = sorted(
            highlights,
            key=lambda h: (h.get("line_number", 1), h.get("start_col", 0)),
            reverse=True
        )

        lines = content.splitlines(keepends=True)
        for h in sorted_highlights:
            line_no = h.get("line_number", 1) - 1
            start = h.get("start_col", 0)
            end = h.get("end_col", 0)
            match_txt = h.get("match_text", "")
            masked = mask_text(match_txt, mask_pattern)

            if 0 <= line_no < len(lines):
                line = lines[line_no]
                if 0 <= start < len(line) and line[start:end] == match_txt:
                    lines[line_no] = line[:start] + masked + line[end:]
                else:
                    # Handle masked match_text
                    if "..." in match_txt and 0 <= start < end <= len(line):
                        seg = line[start:end]
                        if match_txt.startswith(seg[:4]) or match_txt.endswith(seg[-4:]):
                            lines[line_no] = line[:start] + masked + line[end:]
                            continue
                            
                    if match_txt in line:
                        lines[line_no] = line.replace(match_txt, masked, 1)

        updated_content = "".join(lines)
        with open(file_path, "w", encoding=used_encoding) as f:
            f.write(updated_content)

        new_checksum = calculate_file_checksum(file_path)
        _sync_state_after_remediation(file_path, new_checksum)

        return {
            "success": True,
            "file": file_path,
            "redacted_count": len(highlights),
            "new_checksum": new_checksum,
            "backup_path": backup_entry["backup_path"]
        }
    except Exception as e:
        return {"success": False, "error": str(e), "message": f"Error during batch redaction: {str(e)}"}


def _sync_state_after_remediation(file_path: str, new_checksum: str):
    """Updates results and cache after a file has been remediated."""
    from .scanner import get_file_text_content
    new_text = get_file_text_content(Path(file_path))
    remaining_highlights = locate_text_pii_matches(new_text)

    cache = load_cache()
    try:
        st = os.stat(file_path)
        mtime = st.st_mtime
        size = st.st_size
    except Exception:
        mtime = 0
        size = 0

    if not remaining_highlights:
        # Completely clean
        cache[file_path] = {
            "mtime": mtime,
            "size": size,
            "checksum": new_checksum,
            "result": {"compromised": False, "reason": "Remediated by in-place redaction"}
        }
        results = load_results()
        new_results = [r for r in results if r.get("file") != file_path]
        save_results(new_results)

        with state_lock:
            scan_state.flagged_files = [f for f in scan_state.flagged_files if f.get("file") != file_path]
            scan_state.progress["flagged_count"] = len(scan_state.flagged_files)
    else:
        # Update remaining findings
        cache[file_path] = {
            "mtime": mtime,
            "size": size,
            "checksum": new_checksum,
            "result": {
                "compromised": True,
                "reason": f"Flagged with {len(remaining_highlights)} remaining PII findings",
                "snippets": [h["match_text"] for h in remaining_highlights]
            }
        }
        results = load_results()
        for r in results:
            if r.get("file") == file_path:
                r["checksum"] = new_checksum
                r["reason"] = f"Flagged with {len(remaining_highlights)} remaining PII findings"
                r["snippets"] = [h["match_text"] for h in remaining_highlights]
        save_results(results)

        with state_lock:
            for f in scan_state.flagged_files:
                if f.get("file") == file_path:
                    f["checksum"] = new_checksum
                    f["reason"] = f"Flagged with {len(remaining_highlights)} remaining PII findings"
                    f["snippets"] = [h["match_text"] for h in remaining_highlights]

    save_cache(cache)


# --------------------------------------------------------------------------
# System Trash & Deletion Management
# --------------------------------------------------------------------------

def trash_or_delete_file(file_path: str, permanent: bool = False, base_dir: str = ".") -> dict:
    """
    Removes a file from the filesystem.
    By default, moves it to the OS Recycle Bin / System Trash.
    If permanent=True or trash is unavailable, permanently deletes with os.remove.
    """
    if not os.path.exists(file_path):
        _remove_from_state(file_path)
        return {"success": True, "file": file_path, "trashed": False, "permanent": permanent}

    trashed_successfully = False

    if not permanent:
        # Try OS system trash
        if sys.platform == "win32":
            try:
                import ctypes
                from ctypes import wintypes

                class SHFILEOPSTRUCTW(ctypes.Structure):
                    _fields_ = [
                        ("hwnd", wintypes.HWND),
                        ("wFunc", wintypes.UINT),
                        ("pFrom", wintypes.LPCWSTR),
                        ("pTo", wintypes.LPCWSTR),
                        ("fFlags", ctypes.c_uint16),
                        ("fAnyOperationsAborted", wintypes.BOOL),
                        ("hNameMappings", wintypes.LPVOID),
                        ("lpszProgressTitle", wintypes.LPCWSTR)
                    ]

                FO_DELETE = 0x0003
                FOF_ALLOWUNDO = 0x0040
                FOF_NOCONFIRMATION = 0x0010
                FOF_SILENT = 0x0004
                FOF_NOERRORUI = 0x0400

                path_double_null = str(Path(file_path).resolve()) + "\0\0"
                file_op = SHFILEOPSTRUCTW()
                file_op.hwnd = None
                file_op.wFunc = FO_DELETE
                file_op.pFrom = path_double_null
                file_op.pTo = None
                file_op.fFlags = FOF_ALLOWUNDO | FOF_NOCONFIRMATION | FOF_SILENT | FOF_NOERRORUI

                res = ctypes.windll.shell32.SHFileOperationW(ctypes.byref(file_op))
                if res == 0 and not file_op.fAnyOperationsAborted:
                    trashed_successfully = True
            except Exception as e:
                print(f"Windows Recycle Bin error: {e}")

        elif sys.platform == "darwin":
            try:
                import subprocess
                clean_path = str(Path(file_path).resolve()).replace('"', '\\"')
                script = f'tell application "Finder" to delete POSIX file "{clean_path}"'
                res = subprocess.run(["osascript", "-e", script], capture_output=True)
                if res.returncode == 0:
                    trashed_successfully = True
            except Exception as e:
                print(f"macOS Trash error: {e}")

        else:
            # Linux gio / trash-cli fallback
            try:
                import subprocess
                res = subprocess.run(["gio", "trash", str(file_path)], capture_output=True)
                if res.returncode == 0:
                    trashed_successfully = True
            except Exception:
                pass

        # Fallback local trash folder if OS trash didn't work
        if not trashed_successfully and os.path.exists(file_path):
            try:
                trash_dir = os.path.join(base_dir, TRASH_DIR_NAME)
                os.makedirs(trash_dir, exist_ok=True)
                dest = os.path.join(trash_dir, f"{int(time.time())}_{os.path.basename(file_path)}")
                shutil.move(file_path, dest)
                trashed_successfully = True
            except Exception:
                pass

    if not trashed_successfully:
        # Permanent delete fallback
        try:
            if os.path.isdir(file_path):
                shutil.rmtree(file_path)
            else:
                os.remove(file_path)
        except Exception as e:
            return {"success": False, "error": str(e), "file": file_path}

    _remove_from_state(file_path)
    return {
        "success": True,
        "file": file_path,
        "trashed": trashed_successfully,
        "permanent": permanent or not trashed_successfully
    }


def _remove_from_state(file_path: str):
    """Removes a deleted file from scan state and saved results."""
    results = load_results()
    new_results = [r for r in results if r.get("file") != file_path]
    save_results(new_results)

    cache = load_cache()
    if file_path in cache:
        del cache[file_path]
        save_cache(cache)

    with state_lock:
        scan_state.flagged_files = [f for f in scan_state.flagged_files if f.get("file") != file_path]
        scan_state.progress["flagged_count"] = len(scan_state.flagged_files)


# --------------------------------------------------------------------------
# Allowed Exceptions & .argusignore Whitelisting
# --------------------------------------------------------------------------

def load_argusignore(base_dir: str = ".") -> list[str]:
    """Reads lines from .argusignore file."""
    ignore_path = os.path.join(base_dir, ARGUSIGNORE_FILENAME)
    if os.path.exists(ignore_path):
        try:
            with open(ignore_path, "r", encoding="utf-8") as f:
                return [line.strip() for line in f if line.strip() and not line.startswith("#")]
        except Exception:
            pass
    return []


def append_argusignore_entry(rule: str, comment: str = "", base_dir: str = "."):
    """Appends an exception rule to .argusignore."""
    ignore_path = os.path.join(base_dir, ARGUSIGNORE_FILENAME)
    try:
        with open(ignore_path, "a", encoding="utf-8") as f:
            if comment:
                f.write(f"\n# {comment}\n")
            f.write(f"{rule}\n")
    except Exception as e:
        print(f"Error appending to .argusignore: {e}")


def mark_as_safe_exception(
    file_path: str,
    match_text: str = None,
    pattern_name: str = None,
    reason: str = "Whitelisted by user",
    base_dir: str = "."
) -> dict:
    """
    Whitelists a file or specific entity exception.
    Appends to .argusignore, updates settings allowed_exceptions, and clears from results.
    """
    settings = load_settings()
    exceptions = settings.get("allowed_exceptions", [])

    entry_id = str(uuid.uuid4())[:8]
    exception_entry = {
        "id": entry_id,
        "file": os.path.abspath(file_path) if file_path else "",
        "filename": os.path.basename(file_path) if file_path else "",
        "match_text": match_text or "",
        "pattern_name": pattern_name or "",
        "added_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "reason": reason
    }
    exceptions.append(exception_entry)
    settings["allowed_exceptions"] = exceptions
    save_settings(settings)

    # Append to .argusignore
    if file_path:
        norm_path = file_path.replace("\\", "/")
        rule = norm_path
        if match_text:
            rule = f"{norm_path} :: {match_text}"
        append_argusignore_entry(rule, comment=f"Exception {entry_id}: {pattern_name or reason}", base_dir=base_dir)

    # Update cache
    cache = load_cache()
    if file_path:
        try:
            st = os.stat(file_path)
            mtime = st.st_mtime
            size = st.st_size
        except Exception:
            mtime = 0
            size = 0
        checksum = calculate_file_checksum(file_path)
        cache[file_path] = {
            "mtime": mtime,
            "size": size,
            "checksum": checksum,
            "result": {"compromised": False, "marked_ok": True, "reason": reason}
        }
        save_cache(cache)

    _remove_from_state(file_path)
    return {"success": True, "exception": exception_entry}


def get_allowed_exceptions() -> list[dict]:
    """Returns all whitelisted exceptions configured in settings."""
    settings = load_settings()
    return settings.get("allowed_exceptions", [])


def remove_allowed_exception(exception_id: str, base_dir: str = ".") -> dict:
    """
    Removes a whitelisted exception from settings and clears cache for that file.
    """
    settings = load_settings()
    exceptions = settings.get("allowed_exceptions", [])
    removed_entry = None
    new_exceptions = []

    for ex in exceptions:
        if ex.get("id") == exception_id:
            removed_entry = ex
        else:
            new_exceptions.append(ex)

    if not removed_entry:
        return {"success": False, "error": "Exception ID not found"}

    settings["allowed_exceptions"] = new_exceptions
    save_settings(settings)

    # Clear cache entry for the file to re-scan
    file_path = removed_entry.get("file")
    if file_path:
        cache = load_cache()
        if file_path in cache:
            del cache[file_path]
            save_cache(cache)

    return {"success": True, "removed": removed_entry}


def _matches_pattern_or_mask(candidate: str, rule: str) -> bool:
    if not candidate or not rule:
        return False
    if candidate.lower() == rule.lower():
        return True
    if "..." in candidate and len(rule) > 8:
        cand_parts = candidate.split("...")
        if len(cand_parts) == 2 and rule.lower().startswith(cand_parts[0].lower()) and rule.lower().endswith(cand_parts[1].lower()):
            return True
    if "..." in rule and len(candidate) > 8:
        rule_parts = rule.split("...")
        if len(rule_parts) == 2 and candidate.lower().startswith(rule_parts[0].lower()) and candidate.lower().endswith(rule_parts[1].lower()):
            return True
    return False


def is_file_or_match_ignored(file_path: str, match_text: str = None, pattern_name: str = None, base_dir: str = ".") -> bool:
    """
    Determines if a file or finding is whitelisted via .argusignore or settings.
    """
    norm_path = os.path.abspath(file_path).replace("\\", "/").lower()
    
    # 1. Check settings allowed_exceptions
    settings = load_settings()
    exceptions = settings.get("allowed_exceptions", [])
    for ex in exceptions:
        ex_file = ex.get("file", "").replace("\\", "/").lower()
        ex_match = ex.get("match_text", "")
        if ex_file and (ex_file == norm_path or norm_path.endswith(ex_file)):
            if not ex_match:
                return True # Whole file ignored
            if match_text and _matches_pattern_or_mask(match_text, ex_match):
                return True # Specific match ignored

    # 2. Check .argusignore
    import fnmatch
    filename = os.path.basename(norm_path)
    ignore_lines = load_argusignore(base_dir)
    for line in ignore_lines:
        line_clean = line.replace("\\", "/").lower().strip()
        if not line_clean:
            continue
        if "::" in line_clean:
            path_part, match_part = line_clean.split("::", 1)
            path_part = path_part.strip()
            match_part = match_part.strip()
            if (path_part == norm_path or norm_path.endswith(path_part) or fnmatch.fnmatch(filename, path_part) or fnmatch.fnmatch(norm_path, path_part)) and match_text and _matches_pattern_or_mask(match_text, match_part):
                return True
        else:
            if (line_clean == norm_path or 
                norm_path.endswith(line_clean) or 
                filename == line_clean or 
                fnmatch.fnmatch(filename, line_clean) or 
                fnmatch.fnmatch(norm_path, f"*{line_clean.strip('/')}*") or 
                (match_text and _matches_pattern_or_mask(match_text, line_clean))):
                return True

    return False

